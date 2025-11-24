# services/files/document.py
from __future__ import annotations

from typing import Tuple, List
from pathlib import Path

from .types import FileProcessResult
from services.regex_rules import PATTERNS  # 정규식 패턴 전체 사용

# 선택적 의존성 (없으면 해당 포맷만 에러 리턴)
try:
    from docx import Document as DocxDocument  # python-docx
except ImportError:  # pragma: no cover
    DocxDocument = None  # type: ignore

try:
    from pptx import Presentation  # python-pptx
except ImportError:  # pragma: no cover
    Presentation = None  # type: ignore

try:
    from openpyxl import load_workbook  # openpyxl
except ImportError:  # pragma: no cover
    load_workbook = None  # type: ignore


# ---------------------------
# 공통: 정규식 기반 문자열 치환
# ---------------------------

def _mask_text_with_patterns(text: str) -> Tuple[str, bool]:
    """
    regex_rules.PATTERNS 에 있는 모든 패턴을 적용해서
    매칭 문자열을 해당 라벨 토큰(LABEL)로 치환.

    return: (masked_text, changed)
    """
    if not text:
        return text, False

    changed = False
    out = text

    for label, rx in PATTERNS.items():
        token = label  # PHONE, EMAIL, CARD_NUMBER, ...
        new_out, n = rx.subn(token, out)
        if n > 0:
            changed = True
            out = new_out

    return out, changed


def _make_detection_path(path: Path) -> Path:
    """원본 경로에서 *.detection.* 경로 생성."""
    suffix = path.suffix or ""
    return path.with_name(f"{path.stem}.detection{suffix}")


# ---------------------------
# DOCX
# ---------------------------

def _mask_docx_to_detection(path: Path) -> Tuple[bool, str, str]:
    """
    DOCX 내부 문자열을 정규식으로 토큰 치환 후
    원본은 유지하고, *.detection.docx 로 저장.
    """
    if DocxDocument is None:
        return False, "", "python_docx_not_available"

    try:
        doc = DocxDocument(str(path))
        collected: List[str] = []

        # 일반 단락
        for para in doc.paragraphs:
            original = para.text or ""
            masked, changed = _mask_text_with_patterns(original)
            collected.append(masked)
            if changed:
                # 스타일은 일부 깨질 수 있지만 단순하게 전체 텍스트 교체
                para.text = masked

        # 표 안의 셀 텍스트
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    original = cell.text or ""
                    masked, changed = _mask_text_with_patterns(original)
                    collected.append(masked)
                    if changed:
                        cell.text = masked

        out_path = _make_detection_path(path)
        doc.save(str(out_path))
        return True, "\n".join(collected), ""
    except Exception as e:
        return False, "", f"docx_redaction_error:{e}"


# ---------------------------
# PPTX
# ---------------------------

def _mask_pptx_to_detection(path: Path) -> Tuple[bool, str, str]:
    """
    PPTX 슬라이드 내 텍스트를 토큰 치환 후
    원본은 유지하고, *.detection.pptx 로 저장.
    """
    if Presentation is None:
        return False, "", "python_pptx_not_available"

    try:
        prs = Presentation(str(path))
        collected: List[str] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text or ""
                masked, changed = _mask_text_with_patterns(text)
                collected.append(masked)
                if changed:
                    shape.text = masked

        out_path = _make_detection_path(path)
        prs.save(str(out_path))
        return True, "\n".join(collected), ""
    except Exception as e:
        return False, "", f"pptx_redaction_error:{e}"


# ---------------------------
# XLSX
# ---------------------------

def _mask_xlsx_to_detection(path: Path) -> Tuple[bool, str, str]:
    """
    XLSX 각 셀의 문자열 값을 토큰 치환 후
    원본은 유지하고, *.detection.xlsx 로 저장.
    """
    if load_workbook is None:
        return False, "", "openpyxl_not_available"

    try:
        wb = load_workbook(filename=str(path))
        collected: List[str] = []

        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    val = cell.value
                    if isinstance(val, str) and val:
                        masked, changed = _mask_text_with_patterns(val)
                        collected.append(masked)
                        if changed:
                            cell.value = masked

        out_path = _make_detection_path(path)
        wb.save(str(out_path))
        return True, "\n".join(collected), ""
    except Exception as e:
        return False, "", f"xlsx_redaction_error:{e}"


def extract_text_from_office(path: Path) -> Tuple[bool, str, str]:
    """
    DOCX / PPTX / XLSX 텍스트 기반 문서에 대해
    - 파일을 직접 열어 텍스트를 정규식으로 탐지
    - 매칭 구간을 라벨 토큰으로 치환
    - 원본은 유지하고 *.detection.* 파일로 저장
    - 치환된 전체 텍스트를 반환
    """
    suffix = path.suffix.lower()

    if suffix == ".docx":
        return _mask_docx_to_detection(path)
    if suffix == ".pptx":
        return _mask_pptx_to_detection(path)
    if suffix == ".xlsx":
        return _mask_xlsx_to_detection(path)

    return False, "", f"unsupported_office_ext:{suffix}"


# ---------------------------
# TXT / CSV
# ---------------------------

def extract_text_from_plain(path: Path) -> Tuple[bool, str, str]:
    """
    TXT / CSV:
    - 원본 파일을 UTF-8 텍스트로 읽어서 정규식 기반 토큰 치환
    - 원본은 유지하고 *.detection.txt / *.detection.csv 로 저장
    - 치환된 텍스트 반환
    """
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        return False, "", f"plain_text_read_error:{e}"

    masked, _ = _mask_text_with_patterns(raw)
    out_path = _make_detection_path(path)

    try:
        out_path.write_text(masked, encoding="utf-8")
    except Exception as e:
        # 텍스트는 추출했지만 저장 실패
        return True, masked, f"plain_text_write_error:{e}"

    return True, masked, ""


# ---------------------------
# 진입점
# ---------------------------

def process_document_file(result: FileProcessResult) -> FileProcessResult:
    """
    텍스트 기반 문서(docx, pptx, csv, txt, xlsx)에 대한 처리:
    - 원본(saved_path)을 직접 열어 텍스트를 정규식으로 중요정보 탐지
    - 매칭 구간을 라벨 토큰으로 치환
    - 원본은 그대로 두고, *.detection.* 파일로 저장
    - 치환된 전체 텍스트를 extracted_text 에 기록

    PDF는 redaction.py 에서 좌표 기반 레댁션으로 처리한다.
    """
    ext = (result.ext or "").lower()
    path = result.saved_path

    if ext in {"docx", "pptx", "xlsx"}:
        used, text, err = extract_text_from_office(path)
    elif ext in {"txt", "csv"}:
        used, text, err = extract_text_from_plain(path)
    else:
        used, text, err = False, "", f"unsupported_document_ext:{ext}"

    result.ocr_used = used
    result.extracted_text = text if used and text else ""
    result.ocr_error = err or None

    # 나중에 반환할 때 detection 파일을 쓰고 싶다면,
    # 같은 규칙으로 경로를 계산해서 사용하면 됨:
    # detection_path = _make_detection_path(path)
    # (여기서 필드로 붙이고 싶다면 result.detection_path = detection_path 식으로 추가 가능)

    return result
